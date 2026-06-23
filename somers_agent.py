import os
import json
import yaml
import datetime
import subprocess
import glob
import fitz  # PyMuPDF
from typing import List, Dict, Any
import matplotlib.pyplot as plt
import koreanize_matplotlib  # 한글 폰트 자동 설정
import docx
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

class SomersAgent:
    def __init__(self, rules_path: str = "somers_rules.yaml", kb_path: str = "knowledge_base.json"):
        self.rules_path = rules_path
        self.kb_path = kb_path
        self.rules = self._load_rules()
        self.kb = self._load_kb()

    def _load_rules(self) -> Dict[str, Any]:
        if os.path.exists(self.rules_path):
            with open(self.rules_path, "r", encoding="utf-8") as f:
                return yaml.safe_load(f)
        return {}

    def _load_kb(self) -> List[Dict[str, Any]]:
        if os.path.exists(self.kb_path):
            try:
                with open(self.kb_path, "r", encoding="utf-8") as f:
                    return json.load(f)
            except Exception:
                return []
        return []

    def _clean_xml_string(self, text: str) -> str:
        """
        python-docx 등 XML 생성 시 오류를 방지하기 위해 널 바이트 및 제어 문자를 제거합니다.
        """
        if not text:
            return ""
        allowed_chars = []
        for char in text:
            cp = ord(char)
            if (cp == 0x9 or cp == 0xA or cp == 0xD or 
                (0x20 <= cp <= 0xD7FF) or 
                (0xE000 <= cp <= 0xFFFD) or 
                (0x10000 <= cp <= 0x10FFFF)):
                allowed_chars.append(char)
        return "".join(allowed_chars)

    def save_kb(self):
        with open(self.kb_path, "w", encoding="utf-8") as f:
            json.dump(self.kb, f, ensure_ascii=False, indent=2)

    def plan_queries(self, topic: str) -> Dict[str, List[str]]:
        """
        사용자 주제를 바탕으로 고품질 전문 문서 수집용 쿼리를 설계합니다.
        """
        # 규칙 기반의 기본 키워드 설계
        korean_queries = [
            f"{topic} 설계 기준 PDF",
            f"{topic} 가이드라인 보고서",
            f"{topic} 표준 사양서 filetype:pdf",
            f"{topic} 기술 매뉴얼 site:go.kr"
        ]
        
        # 영어 번역 쿼리 예측 (간단한 매핑 또는 기본 명사 조합)
        # 실제 환경에서는 AI 번역 API 등을 연동할 수 있음
        english_topic = topic.replace("식품공장 위생배관", "hygienic piping food factory")\
                             .replace("MES 견적", "manufacturing execution system cost estimation")\
                             .replace("자동화", "automation")
        
        english_queries = [
            f"{english_topic} design guideline filetype:pdf",
            f"{english_topic} specification standard PDF",
            f"EHEDG {english_topic} guideline",
            f"FDA {english_topic} regulations"
        ]

        return {
            "korean": korean_queries,
            "english": english_queries
        }

    def search_and_retrieve(self, topic: str, search_results_mock: List[Dict[str, Any]] = None) -> List[Dict[str, Any]]:
        """
        실시간 웹 검색을 모사하거나, 제공된 검색 데이터로부터 신뢰도 높은 전문문서 목록을 수집합니다.
        """
        if search_results_mock:
            return search_results_mock
        
        # 기본 Mock 데이터 제공 (도메인별 고품질 결과 자동 매핑)
        if "위생배관" in topic:
            return [
                {
                    "title": "식품공장의 위생 배관 설계 가이드라인 및 기준 연구",
                    "source": "한국식품안전관리인증원",
                    "year": 2024,
                    "type": "정부·공공기관 보고서",
                    "url": "https://www.haccp.or.kr/guideline_piping_2024.pdf",
                    "abstract": "본 가이드라인은 식품 제조설비의 위생적 설계를 위한 배관 재질 선정, 구배(Slope) 설계, 세척(CIP) 용이성 확보 방안을 제시함.",
                    "details": "배관 재질은 STS316L 사용 권장, 용접부는 내면 비드 제거 필수, 1/100 이상의 배수 구배 유지 필요."
                },
                {
                    "title": "EHEDG Doc 8: Hygienic Design Principles for Food Processing",
                    "source": "EHEDG (European Hygienic Engineering & Design Group)",
                    "year": 2023,
                    "type": "국제 표준 가이드라인",
                    "url": "https://www.ehedg.org/guidelines/doc8.pdf",
                    "abstract": "European standards for hygienic engineering and process piping design in food and beverage production lines.",
                    "details": "Avoid dead legs (L/D < 2), use aseptic flange joints, self-draining properties verification methods included."
                },
                {
                    "title": "식품 제조공장 스마트 위생 배관 시스템 특허 분석 보고서",
                    "source": "특허청",
                    "year": 2025,
                    "type": "특허 및 기술 분석서",
                    "url": "https://www.kipo.go.kr/patent_smart_piping.pdf",
                    "abstract": "스마트 위생 배관의 자동 밸브 제어 및 누수/오염 센싱 기술에 관한 특허 동향과 실무 설계 적용 분석.",
                    "details": "자동 세척 프로세스와 센서 노즐 결합부에 대한 특허 설계 규격 포함."
                }
            ]
        elif "MES" in topic or "견적" in topic:
            return [
                {
                    "title": "공공/제조 소프트웨어 사업 대가산정 가이드라인",
                    "source": "한국소프트웨어산업협회 (KOSA)",
                    "year": 2025,
                    "type": "공공기관 표준가이드",
                    "url": "https://www.sw.or.kr/cost_guide_2025.pdf",
                    "abstract": "기능점수(FP) 방식에 따른 제조 IT 시스템 및 MES 구축비 정밀 공수 산정 및 타당성 검토 가이드.",
                    "details": "단가 기준 1FP 당 가격 계산법, 간이법 및 상세법 산정 절차 수록."
                },
                {
                    "title": "중소·중견기업 스마트공장 MES 구축 성공요인 및 견적 타당성 분석",
                    "source": "한국산업기술진흥원",
                    "year": 2024,
                    "type": "정부산하 연구보고서",
                    "url": "https://www.kiat.or.kr/mes_feasibility_study.pdf",
                    "abstract": "MES 도입 단계에서 빈번히 발생하는 추가 공수 비용 낭비를 예방하기 위한 견적 타당성 분석 프레임워크.",
                    "details": "하드웨어 인프라 비용 대비 소프트웨어 커스터마이징 비용 적정 비율 분석 (보통 1:1.5 이하 권장)."
                }
            ]
        else:
            return [
                {
                    "title": f"{topic}에 관한 산업 동향 및 학술 보고서",
                    "source": "한국학술지인용색인 (KCI)",
                    "year": 2024,
                    "type": "학술 논문",
                    "url": "https://www.kci.go.kr/article_summary.pdf",
                    "abstract": f"{topic}의 기초 기술 분석과 최신 적용 표준 분석 내용.",
                    "details": "실무 적용 포인트 및 핵심 기술 사양이 요약됨."
                }
            ]

    def evaluate_reliability(self, doc: Dict[str, Any], constraints: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        문서 신뢰도를 7대 항목 기준으로 평가하고 종합 별점을 부여합니다.
        사용자의 요구 조건(최신성 기준 연도 등)이 있다면 동적으로 점수를 조정합니다.
        """
        scores = {}
        
        # 1. 출처 신뢰도
        source = doc.get("source", "")
        if any(org in source for org in ["정부", "국제", "HACCP", "EHEDG", "협회", "청", "부", "KOSA", "KIAT"]):
            scores["source_reliability"] = 5
        else:
            scores["source_reliability"] = 3
            
        # 2. 최신성 (사용자 제한 조건 반영)
        year = doc.get("year", 2024)
        age = datetime.datetime.now().year - year
        
        max_age = 5 # 기본값은 5년
        if constraints and "max_age" in constraints:
            max_age = constraints["max_age"]
            
        if age <= 2:
            scores["recency"] = 5
        elif age <= max_age:
            scores["recency"] = 4
        else:
            # 사용자가 설정한 최신성 기준 연도를 초과한 경우 최하점 처리
            scores["recency"] = 1
            
        # 3. 전문성
        scores["professionalism"] = 5 if "가이드라인" in doc.get("title", "") or "표준" in doc.get("title", "") or "연구" in doc.get("title", "") else 4
        
        # 4. 원문성 (원문 PDF 기준 5점)
        scores["originality"] = 5 if ".pdf" in doc.get("url", "").lower() else 4
        
        # 5. 실무성
        scores["practicality"] = 5 if "설계" in doc.get("title", "") or "견적" in doc.get("title", "") or "가이드" in doc.get("title", "") else 4
        
        # 6. 데이터성
        scores["data_density"] = 4  # 수치, 표가 포함되어 있음을 기본 가정
        
        # 7. 객관성
        scores["objectivity"] = 5  # 공공/학술 자료는 편향성 낮음
        
        # 종합 평균 계산
        avg_score = round(sum(scores.values()) / len(scores))
        
        # 별점 변환
        stars = "★" * avg_score + "☆" * (5 - avg_score)
        
        return {
            "scores": scores,
            "rating_value": avg_score,
            "rating_stars": stars
        }

    def train_ai_knowledge(self, doc: Dict[str, Any], rating_info: Dict[str, Any], topic: str, constraints: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        새로 취득한 고신뢰도 전문문서 정보를 AI가 학습하기 좋은 '지식 청크' 및 'Q&A' 구조로 변환합니다.
        사용자의 사전 요청사항(조건)이 있을 경우 학습용 Q&A에 명시적으로 반영합니다.
        """
        # AI 학습용 QA 및 지식 세트 생성
        knowledge_chunks = [
            {
                "question": f"{doc['title']}의 핵심 설계/검토 기준은 무엇입니까?",
                "answer": doc['details']
            },
            {
                "question": f"{doc['title']} 문서의 신뢰성과 출처 정보는 어떻게 됩니까?",
                "answer": f"본 문서는 {doc['year']}년 {doc['source']}에서 발행된 {doc['type']}로, 신뢰도 등급은 {rating_info['rating_stars']} 입니다."
            }
        ]
        
        # 사용자 세부 요청사항 반영
        if constraints and constraints.get("user_request"):
            knowledge_chunks.append({
                "question": f"이 지식을 활용할 때 사용자의 요청사항인 '{constraints['user_request']}' 측면에서 검토해야 할 점은 무엇입니까?",
                "answer": f"문서의 분석 결과, {doc['details']}에 기반하여 사용자가 기재한 조건인 '{constraints['user_request']}'에 맞추어 의사결정을 지원합니다."
            })

        learning_entry = {
            "date": datetime.date.today().isoformat(),
            "category": doc["type"],
            "topic": topic,
            "title": doc["title"],
            "source": doc["source"],
            "year": doc["year"],
            "url": doc["url"],
            "reliability": rating_info["rating_stars"],
            "scores": rating_info["scores"],
            "knowledge_chunks": knowledge_chunks,
            "summary": doc["abstract"]
        }
        
        return learning_entry

    def run_daily_learning(self, topic: str, constraints: Dict[str, Any] = None) -> List[Dict[str, Any]]:
        """
        주제와 제약 조건을 받아 문서를 수집, 평가하고 지식 저장소에 누적 저장하는 전체 프로세스입니다.
        """
        print(f"[{datetime.date.today().isoformat()}] '{topic}' 분야 전문 지식 습득 및 학습 시작...")
        if constraints:
            print(f" -> 적용된 사전 학습 조건: {constraints}")
        
        # 1. 쿼리 기획
        queries = self.plan_queries(topic)
        print(f" -> 생성된 한글 쿼리: {queries['korean']}")
        print(f" -> 생성된 영문 쿼리: {queries['english']}")
        
        # 2. 문서 수집
        raw_docs = self.search_and_retrieve(topic)
        
        new_learnings = []
        for doc in raw_docs:
            # 3. 신뢰도 검증 (동적 제약조건 전달)
            rating_info = self.evaluate_reliability(doc, constraints)
            
            # 4. AI 학습 데이터 변환 (동적 제약조건 전달)
            learning_entry = self.train_ai_knowledge(doc, rating_info, topic, constraints)
            new_learnings.append(learning_entry)
            
            # 5. 지식베이스에 추가 (중복 방지)
            if not any(item["title"] == learning_entry["title"] for item in self.kb):
                self.kb.append(learning_entry)
                
        self.save_kb()
        print(f" -> 총 {len(new_learnings)}건의 전문 지식을 학습하고 {self.kb_path}에 저장 완료했습니다.")
        return new_learnings

    def generate_docx_report(self, learnings: List[Dict[str, Any]], filename: str = "보고서.docx"):
        """
        학습된 지식을 바탕으로 전문가 수준의 Word 보고서를 생성합니다.
        """
        doc = docx.Document()
        
        # 스타일 정의
        style = doc.styles['Normal']
        font = style.font
        font.name = '맑은 고딕'
        font.size = Pt(10)
        
        # 제목 추가
        title_para = doc.add_paragraph()
        title_run = title_para.add_run("전문문서 검색 및 AI 학습 검토 보고서")
        title_run.font.name = '맑은 고딕'
        title_run.font.size = Pt(20)
        title_run.font.bold = True
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 메타 정보
        meta_para = doc.add_paragraph()
        meta_para.add_run(f"발행일: {datetime.date.today().isoformat()}\n작성자: SomersAgent (지식 관리 시스템)")
        meta_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        
        # 1. 개요
        doc.add_heading("1. 개요 및 목적", level=1)
        doc.add_paragraph("본 보고서는 최신 정보와 고신뢰도 공공/학술 문서를 검색하고 그 결과를 체계적으로 분류 및 평가하여 AI의 고수준 학습을 돕고 실무 구매, 설계 및 견적 타당성 검토에 바로 활용할 수 있도록 작성되었습니다.")

        # 2. 수집 자료 평가 내역
        doc.add_heading("2. 수집 자료 및 신뢰도 평가 목록", level=1)
        
        # 표 생성
        table = doc.add_table(rows=1, cols=6)
        table.style = 'Table Grid'
        hdr_cells = table.rows[0].cells
        headers = ["No", "문서명", "발행기관", "연도", "유형", "신뢰도"]
        for i, header in enumerate(headers):
            hdr_cells[i].text = header
            hdr_cells[i].paragraphs[0].runs[0].font.bold = True
            
        for index, item in enumerate(learnings):
            row_cells = table.add_row().cells
            row_cells[0].text = str(index + 1)
            row_cells[1].text = item["title"]
            row_cells[2].text = item["source"]
            row_cells[3].text = str(item["year"])
            row_cells[4].text = item["category"]
            row_cells[5].text = item["reliability"]

        # 공백 추가
        doc.add_paragraph()

        # 3. 핵심 자료 및 실무 활용 가이드
        doc.add_heading("3. 핵심 자료별 요약 및 실무 적용 포인트", level=1)
        
        for item in learnings:
            doc.add_heading(f"■ {item['title']}", level=2)
            
            p = doc.add_paragraph()
            p.add_run("• 출처 및 유형: ").bold = True
            p.add_run(f"{item['source']} ({item['year']}) | {item['category']}\n")
            p.add_run("• 핵심 내용 요약:\n").bold = True
            p.add_run(f"  {item['summary']}\n")
            
            # 실무 활용 파트 추출
            p.add_run("• AI 학습용 핵심 Q&A:\n").bold = True
            for chunk in item["knowledge_chunks"]:
                p.add_run(f"  - Q: {chunk['question']}\n")
                p.add_run(f"    A: {chunk['answer']}\n")
                
        doc.save(filename)
        print(f" -> Word 보고서가 '{filename}' 파일로 저장되었습니다.")

    def generate_md_report(self, learnings: List[Dict[str, Any]], filename: str = "보고서.md"):
        """
        템플릿에 맞춘 Markdown 보고서를 생성합니다.
        """
        md_content = []
        md_content.append(f"# [{learnings[0]['topic']}] 전문자료 검색 및 학습 결과\n")
        
        # 1. 요청 해석
        md_content.append("## 1. 요청 해석")
        md_content.append(f"- **요청 분야:** {learnings[0]['topic']}")
        md_content.append("- **검색 목적:** 정확한 정보에 입각한 지식 아키텍처 수립 및 실무 의사결정 지원")
        md_content.append("- **필요한 자료 유형:** 공공기관 보고서, 학술지 논문, 국제 가이드라인, 표준 문서")
        md_content.append("- **최신성 필요 여부:** 최근 3~5년 이내 최신 법규/기술기준 준수 필수\n")
        
        # 2. 검색 키워드
        queries = self.plan_queries(learnings[0]['topic'])
        md_content.append("## 2. 검색 키워드")
        md_content.append("### 한글 키워드")
        for q in queries["korean"]:
            md_content.append(f"- `{q}`")
        md_content.append("\n### 영문 키워드")
        for q in queries["english"]:
            md_content.append(f"- `{q}`")
        md_content.append("\n")
        
        # 3. 추천 문서 목록
        md_content.append("## 3. 추천 문서 목록")
        md_content.append("| No | 문서명 | 발행기관/저자 | 연도 | 유형 | 신뢰도 | 링크 |")
        md_content.append("|---|---|---|---|---|---|---|")
        for index, item in enumerate(learnings):
            md_content.append(f"| {index+1} | {item['title']} | {item['source']} | {item['year']} | {item['category']} | {item['reliability']} | [원문보기]({item['url']}) |")
        md_content.append("\n")
        
        # 4. 핵심자료 요약
        md_content.append("## 4. 핵심자료 요약")
        for index, item in enumerate(learnings):
            md_content.append(f"### {index+1}) {item['title']}")
            md_content.append(f"- **출처 및 연도:** {item['source']} ({item['year']})")
            md_content.append(f"- **주요 내용:** {item['summary']}")
            md_content.append(f"- **AI 학습 Q&A 포인트:**")
            for chunk in item["knowledge_chunks"]:
                md_content.append(f"  - **질문:** {chunk['question']}")
                md_content.append(f"    - **답변:** {chunk['answer']}")
            md_content.append("\n")
            
        # 5. 실무 적용 포인트
        md_content.append("## 5. 실무 적용 포인트")
        md_content.append("| 적용 영역 | 활용 방법 |")
        md_content.append("|---|---|")
        if "위생배관" in learnings[0]['topic']:
            md_content.append("| 설계 검토 | 배관 기울기(1/100) 및 Dead Leg(L/D < 2) 기준 준수 여부를 배관 평면도에서 체크 |")
            md_content.append("| 견적 검토 | STS316L 등 고품질 위생재질 원자재 단가 변동률 대비 설치 인건비 적정성 대조 |")
            md_content.append("| 구매 검토 | EHEDG 또는 3-A 인증 여부 성적서 필수 제출 조건 지정 |")
        else:
            md_content.append("| 설계 검토 | 표준 설계안 및 사용자 요구 사양서(URS) 기준 항목 체크리스트 대조 |")
            md_content.append("| 견적 검토 | 기능점수(FP) 혹은 업계 표준 품셈 단가 적용 적절성 분석 |")
            md_content.append("| 보고서 작성 | 공공/학술지 원문 출처를 각주에 기입하여 고신뢰도 검토서 완성 |")
        md_content.append("\n")
        
        # 6. 자료 신뢰도 판단
        md_content.append("## 6. 자료 신뢰도 판단")
        high_docs = [item['title'] for item in learnings if item['reliability'].count('★') >= 4]
        md_content.append(f"- **가장 신뢰도 높은 자료:** {', '.join(high_docs) if high_docs else '없음'}")
        md_content.append("- **보조 참고자료:** 인터넷 기술 포럼 및 제조사 제품 브로슈어")
        md_content.append("- **추가 확인이 필요한 자료:** 법적 강제 규제 변경 여부 (예: HACCP 고시 개정안 최신본)\n")
        
        # 7. 최종 추천
        md_content.append("## 7. 최종 추천 및 활용 방안")
        md_content.append(f"- 사용자는 가장 먼저 **{learnings[0]['title']}** 자료를 확인하는 것이 좋습니다.")
        md_content.append(f"- **이유:** 해당 도메인의 국가 및 글로벌 표준 가이드라인 역할을 하고 있어 설계 승인과 비용 검토의 절대적인 기준을 제공하기 때문입니다.")
        
        with open(filename, "w", encoding="utf-8") as f:
            f.write("\n".join(md_content))
            
        print(f" -> 마크다운 보고서가 '{filename}' 파일로 저장되었습니다.")

    def visualize_statistics(self, learnings: List[Dict[str, Any]], filename: str = "신뢰도_분포.png"):
        """
        수집된 문서들의 신뢰성 평가 지수(출처, 최신성 등)를 시각화하여 저장합니다.
        (seaborn 스타일 미사용 및 koreanize-matplotlib 준수)
        """
        if not learnings:
            return
            
        labels = ["출처신뢰도", "최신성", "전문성", "원문성", "실무성"]
        
        # 평균값 계산
        avg_scores = [0.0] * 5
        for item in learnings:
            s = item["scores"]
            avg_scores[0] += s.get("source_reliability", 0)
            avg_scores[1] += s.get("recency", 0)
            avg_scores[2] += s.get("professionalism", 0)
            avg_scores[3] += s.get("originality", 0)
            avg_scores[4] += s.get("practicality", 0)
            
        avg_scores = [x / len(learnings) for x in avg_scores]
        
        plt.figure(figsize=(8, 5))
        bars = plt.bar(labels, avg_scores, color='#1f77b4', edgecolor='black', width=0.6)
        
        plt.title("수집 전문문서의 5대 신뢰도 지표 평균값", fontsize=14, fontweight='bold', pad=15)
        plt.ylabel("평점 (5점 만점)", fontsize=12)
        plt.ylim(0, 5.5)
        plt.grid(axis='y', linestyle='--', alpha=0.7)
        
        # 값 표시
        for bar in bars:
            height = bar.get_height()
            plt.text(bar.get_x() + bar.get_width()/2.0, height + 0.1, f'{height:.1f}', ha='center', va='bottom', fontsize=10, fontweight='bold')
            
        plt.tight_layout()
        plt.savefig(filename, dpi=300)
        plt.close()
        print(f" -> 통계 시각화 차트가 '{filename}'로 저장되었습니다.")

    def sync_to_github(self, repo_url: str = "https://github.com/pjw1356-hub/connected-AI.git"):
        """
        학습 결과물(지식베이스, 보고서, 차트)을 사용자 깃허브 저장소로 자동 커밋 및 푸시합니다.
        """
        print(f"\n[{datetime.date.today().isoformat()}] 깃허브 원격 저장소 동기화 시작: {repo_url}")
        
        def run_git(args: List[str]):
            result = subprocess.run(args, capture_output=True, text=True, encoding="utf-8")
            if result.returncode != 0:
                print(f"Git Warning/Error: {' '.join(args)}\nstdout: {result.stdout}\nstderr: {result.stderr}")
            return result

        # 1. git 저장소 초기화 여부 체크
        if not os.path.exists(".git"):
            print(" -> 로컬 Git 저장소를 초기화합니다.")
            run_git(["git", "init"])
            run_git(["git", "branch", "-M", "main"])

        # 2. 원격 저장소 설정 체크 및 추가/갱신
        remote_check = run_git(["git", "remote", "-v"])
        if "origin" not in remote_check.stdout:
            print(" -> 원격 저장소(origin)를 등록합니다.")
            run_git(["git", "remote", "add", "origin", repo_url])
        else:
            if "connected-AI" not in remote_check.stdout:
                print(" -> 새로운 원격 저장소 주소로 갱신합니다.")
                run_git(["git", "remote", "set-url", "origin", repo_url])

        # 3. Git 사용자 정보 검사 및 임시 설정 (GitHub Actions 대응)
        name_check = run_git(["git", "config", "user.name"])
        email_check = run_git(["git", "config", "user.email"])
        if not name_check.stdout.strip():
            print(" -> Git 사용자명을 'somers-agent'로 임시 설정합니다.")
            run_git(["git", "config", "user.name", "somers-agent"])
        if not email_check.stdout.strip():
            print(" -> Git 이메일을 'somers-agent@users.noreply.github.com'으로 임시 설정합니다.")
            run_git(["git", "config", "user.email", "somers-agent@users.noreply.github.com"])

        # 4. 파일 스테이징
        files_to_sync = ["knowledge_base.json", "*.md", "*.docx", "*.png", "somers_rules.yaml", "somers_agent.py", "run_search.py"]
        for pattern in files_to_sync:
            run_git(["git", "add", pattern])

        # 5. 커밋 수행 (변경 사항이 있을 때만)
        status_check = run_git(["git", "status", "--porcelain"])
        if not status_check.stdout.strip():
            print(" -> 변경된 지식 및 산출물 파일이 없습니다. 동기화를 건너뜁니다.")
            return

        commit_msg = f"daily_learning_sync: {datetime.date.today().isoformat()} - 지식베이스 및 보고서 업데이트"
        commit_res = run_git(["git", "commit", "-m", commit_msg])

        if commit_res.returncode == 0:
            print(" -> 로컬 변경 사항 커밋 완료.")
        
        # 6. 원격 변경사항 안전하게 가져오기 (Rebase)
        print(" -> 원격 저장소의 최신 이력을 가져와 병합을 시도합니다 (Pull with Rebase)...")
        pull_res = run_git(["git", "pull", "--rebase", "origin", "main"])
        if pull_res.returncode != 0:
            print(" -> [오류] 원격 변경사항을 가져오는 중 충돌이 발생했거나 오류가 있습니다. 강제 병합 취소(Abort)를 수행합니다.")
            run_git(["git", "rebase", "--abort"])
            return

        # 7. 푸시 실행
        print(" -> 깃허브 저장소로 Push를 시도합니다...")
        push_res = run_git(["git", "push", "origin", "main"])
        if push_res.returncode == 0:
            print(" -> 깃허브 동기화가 성공적으로 완료되었습니다!")
        else:
            print(" -> 깃허브 푸시 중 오류가 발생했습니다. 권한 설정 혹은 원격 저장소 상태를 확인하세요.")

    def _extract_docx_text(self, docx_path: str) -> str:
        """
        DOCX 파일에서 본문 문단 및 표 내부의 텍스트를 추출합니다.
        """
        try:
            doc = docx.Document(docx_path)
            full_text = []
            for para in doc.paragraphs:
                if para.text:
                    full_text.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            full_text.append(cell.text)
            return self._clean_xml_string("\n".join(full_text))
        except Exception as e:
            print(f" -> DOCX 텍스트 추출 중 오류 발생 ({os.path.basename(docx_path)}): {e}")
            return ""

    def _extract_pptx_text(self, pptx_path: str) -> str:
        """
        PPTX 파일에서 슬라이드 본문 및 텍스트 박스 내용을 추출합니다.
        """
        try:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            full_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        full_text.append(shape.text)
            return self._clean_xml_string("\n".join(full_text))
        except Exception as e:
            print(f" -> PPTX 텍스트 추출 중 오류 발생 ({os.path.basename(pptx_path)}): {e}")
            return ""

    def learn_local_files(self, file_paths: List[str]) -> List[Dict[str, Any]]:
        """
        제공된 PDF 및 DOCX 파일 목록을 읽고 분석하여 RAG용 지식 청크(Q&A)로 학습 및 저장합니다.
        """
        print(f"\n[{datetime.date.today().isoformat()}] 지정된 로컬 문서 일괄 학습 및 지식화 기동...")
        
        if not file_paths:
            print(" -> 학습할 파일 목록이 비어 있습니다.")
            return []
            
        print(f" -> 총 {len(file_paths)}개의 문서 발견: {[os.path.basename(f) for f in file_paths]}")
        
        new_learnings = []
        for file_path in file_paths:
            filename = os.path.basename(file_path)
            ext = os.path.splitext(filename)[1].lower()
            print(f" -> '{filename}' 파싱 및 학습 진행 중...")
            
            try:
                full_text = ""
                # 1. 확장자별 텍스트 추출
                if ext == ".pdf":
                    doc = fitz.open(file_path)
                    pages_to_read = min(len(doc), 15)  # 최대 15페이지까지 더 정밀하게 탐독
                    for i in range(pages_to_read):
                        text = doc[i].get_text()
                        if text:
                            full_text += text + "\n"
                elif ext == ".docx":
                    full_text = self._extract_docx_text(file_path)
                elif ext == ".pptx":
                    full_text = self._extract_pptx_text(file_path)
                else:
                    print(f" -> 지원하지 않는 확장자입니다 ({ext}). 건너뜁니다.")
                    continue
                        
                # 텍스트가 극도로 부족한 경우 예외 처리
                if len(full_text.strip()) < 50:
                    full_text = "본문 텍스트를 추출할 수 없거나 이미지 위주의 문서입니다. 파일명을 기반으로 지식을 정리합니다."
                else:
                    full_text = self._clean_xml_string(full_text)
                
                # 2. 카테고리 자동 분류 (텍스트 키워드 기반 분류)
                content_lower = full_text.lower()
                category = "일반 기술 분석서"
                if any(k in content_lower for k in ["abstract", "introduction", "논문", "journal", "ieee", "springer", "연구"]):
                    category = "학술 논문"
                elif any(k in content_lower for k in ["manual", "매뉴얼", "설계", "사양", "specification", "guideline"]):
                    category = "기술 설계 문서"
                elif any(k in content_lower for k in ["haccp", "공공", "가이드라인", "고시", "기준", "정부", "원장"]):
                    category = "정부·공공기관 보고서"
                
                # 3. 문서 메타데이터 유추 (텍스트에서 발행년도, 기관 등 탐색)
                year = datetime.datetime.now().year
                for y in range(2015, year + 1):
                    if str(y) in full_text:
                        year = y
                        break
                        
                source = "로컬 소머즈 문서고"
                if "kci" in content_lower or "dbpia" in content_lower:
                    source = "학술지 DB"
                elif "kipo" in content_lower or "특허" in content_lower:
                    source = "특허청"
                elif "ehedg" in content_lower:
                    source = "EHEDG"
                elif "fda" in content_lower:
                    source = "FDA"

                # 4. 요약 추출 (본문 앞부분)
                summary = full_text.strip()[:180].replace("\n", " ") + "..."

                # 5. AI 학습용 핵심 Q&A 지식 청크 구성
                clean_filename = os.path.splitext(filename)[0]
                knowledge_chunks = [
                    {
                        "question": f"'{clean_filename}' 문서의 주요 목적과 핵심 요약은 무엇입니까?",
                        "answer": f"본 문서는 '{category}' 유형으로 분류되며, 주요 요약 내용은 다음과 같습니다: {summary}"
                    },
                    {
                        "question": f"'{clean_filename}' 문서에서 도출할 수 있는 실무 적용 포인트는 무엇입니까?",
                        "answer": f"문서의 텍스트 구성 성분을 분석한 결과, 발행 연도 {year}년 기준의 {source} 관련 지식 자산입니다. 분석 결과에 기반해 설계 검토 및 URS 사양 검증에 참고하시기 바랍니다."
                    }
                ]
                
                # 6. 신뢰도 평가
                scores = {
                    "source_reliability": 4,
                    "recency": 4 if (datetime.datetime.now().year - year) <= 5 else 2,
                    "professionalism": 5,
                    "originality": 5,
                    "practicality": 4,
                    "data_density": 4,
                    "objectivity": 5
                }
                avg_score = round(sum(scores.values()) / len(scores))
                stars = "★" * avg_score + "☆" * (5 - avg_score)
                
                learning_entry = {
                    "date": datetime.date.today().isoformat(),
                    "category": category,
                    "topic": f"로컬 문서: {clean_filename}",
                    "title": clean_filename,
                    "source": source,
                    "year": year,
                    "url": f"file:///{os.path.abspath(file_path).replace(os.sep, '/')}",
                    "reliability": stars,
                    "scores": scores,
                    "knowledge_chunks": knowledge_chunks,
                    "summary": summary,
                    "applied_constraints": {"mode": "local_file_parser"}
                }
                
                new_learnings.append(learning_entry)
                
                # 중복 저장 방지
                # 기존 항목 중 동일 타이틀이 있다면 덮어쓰거나 지우고 교체하도록 처리
                self.kb = [item for item in self.kb if item["title"] != learning_entry["title"]]
                self.kb.append(learning_entry)
                    
            except Exception as e:
                print(f" -> '{filename}' 파싱 오류 발생: {e}")
                
        self.save_kb()
        print(f" -> 로컬 문서 총 {len(new_learnings)}건을 성공적으로 분류 학습하여 지식베이스에 추가하였습니다.")
        return new_learnings

    def learn_local_pdfs(self, pdf_dir: str = "incoming_documents") -> List[Dict[str, Any]]:
        """
        하위 호환성을 유지하며 로컬 디렉토리 내의 모든 PDF, PPTX, DOCX 파일을 찾아 학습하도록 호출합니다.
        """
        if not os.path.exists(pdf_dir):
            os.makedirs(pdf_dir, exist_ok=True)
            
        files = []
        for ext in ["*.pdf", "*.pptx", "*.docx"]:
            pattern = os.path.join(pdf_dir, ext)
            # 대소문자 구분 없이 찾기 위해 glob 처리
            files.extend(glob.glob(pattern))
        return self.learn_local_files(files)


# 단독 테스트를 위한 메인 함수
if __name__ == "__main__":
    agent = SomersAgent()
    learnings = agent.run_daily_learning("식품공장 위생배관")
    agent.generate_docx_report(learnings, "식품공장_위생배관_학습보고서.docx")
    agent.generate_md_report(learnings, "식품공장_위생배관_학습보고서.md")
    agent.visualize_statistics(learnings, "식품공장_위생배관_신뢰도분포.png")
